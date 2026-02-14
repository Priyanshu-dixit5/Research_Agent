"""
ppt_generator.py – PowerPoint Presentation Generator for ScholarMind
Uses python-pptx to build a professional, interactive research
presentation from slide data (JSON) produced by the summarizer.
Features gradient backgrounds, accent shapes, and slide transitions.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ── Multilingual Font Selection ─────────────────────────────────────
# Nirmala UI ships with Windows 10/11 and supports all Indic scripts
FONT_DEFAULT = "Calibri"
FONT_UNICODE = "Nirmala UI"  # Hindi, Marathi, Sanskrit, Tamil, Telugu, etc.
ENGLISH_ONLY = {"English"}

def _get_font(language: str = "English") -> str:
    """Return the font family name to use for the given language."""
    if language in ENGLISH_ONLY:
        return FONT_DEFAULT
    return FONT_UNICODE


# ── Output directory ────────────────────────────────────────────────
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── Color Palettes ──────────────────────────────────────────────────
THEMES = {
    "default": {
        "bg": RGBColor(0x0B, 0x0B, 0x1E),
        "bg_alt": RGBColor(0x0F, 0x0F, 0x28),
        "accent": RGBColor(0x6C, 0x63, 0xFF),
        "accent_light": RGBColor(0x8B, 0x83, 0xFF),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0xB0, 0xB0, 0xCC),
        "bullet": RGBColor(0xE0, 0xE0, 0xF0),
        "slide_title": RGBColor(0x8B, 0x83, 0xFF),
        "number_bg": RGBColor(0x14, 0x14, 0x30),
        "divider": RGBColor(0x6C, 0x63, 0xFF),
    },
}

THEME = THEMES["default"]


def _set_slide_bg(slide, color=None):
    """Set the slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or THEME["bg"]


def _add_gradient_bar(slide, top=False, bottom=False, **kwargs):
    """Add a gradient accent bar at top and/or bottom."""
    if top:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0), top=Inches(0),
            width=Inches(10), height=Inches(0.07),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME["accent"]
        shape.line.fill.background()

    if bottom:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0), top=Inches(7.4),
            width=Inches(10), height=Inches(0.1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME["accent"]
        shape.line.fill.background()
        shape.fill.fore_color.rgb = RGBColor(0x4C, 0x43, 0xCF)


def _add_decorative_circle(slide, left, top, size, opacity_color):
    """Add a decorative circle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=Inches(left), top=Inches(top),
        width=Inches(size), height=Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = opacity_color
    shape.line.fill.background()


def _add_slide_number(slide, number: int, total: int, language: str = "English"):
    """Add a styled slide number indicator at the bottom-right."""
    font_name = _get_font(language)
    # Background pill
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(8.2), top=Inches(6.85),
        width=Inches(1.5), height=Inches(0.4),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = THEME["number_bg"]
    pill.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=Inches(8.2), top=Inches(6.88),
        width=Inches(1.5), height=Inches(0.35),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}  /  {total}"
    p.font.size = Pt(10)
    p.font.name = font_name
    p.font.color.rgb = THEME["subtitle"]
    p.alignment = PP_ALIGN.CENTER


def _add_branding(slide, language: str = "English"):
    """Add ScholarMind branding to bottom-left."""
    font_name = _get_font(language)
    txBox = slide.shapes.add_textbox(
        left=Inches(0.4), top=Inches(6.9),
        width=Inches(2.5), height=Inches(0.35),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ScholarMind"
    p.font.size = Pt(9)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x77)
    p.alignment = PP_ALIGN.LEFT


def _create_title_slide(prs: Presentation, topic: str, bullets: list[str], language: str = "English"):
    """Create the title slide with topic name and subtitle."""
    font_name = _get_font(language)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Top accent bar
    _add_gradient_bar(slide, top=True, bottom=True)

    # Decorative circles
    _add_decorative_circle(slide, 7.5, 1.0, 2.5, RGBColor(0x18, 0x18, 0x38))
    _add_decorative_circle(slide, 8.0, 1.5, 1.5, RGBColor(0x1F, 0x1F, 0x42))
    _add_decorative_circle(slide, -0.5, 5.0, 2.0, RGBColor(0x15, 0x15, 0x30))

    # Left accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(2.4),
        width=Inches(0.12), height=Inches(2.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["accent"]
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(
        left=Inches(1.3), top=Inches(2.0),
        width=Inches(7.5), height=Inches(1.8),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = THEME["title"]
    p.alignment = PP_ALIGN.LEFT

    # Subtitle lines
    y_offset = 4.2
    for i, bullet_text in enumerate(bullets[:3]):
        txSub = slide.shapes.add_textbox(
            left=Inches(1.3), top=Inches(y_offset + i * 0.5),
            width=Inches(7), height=Inches(0.45),
        )
        tf_sub = txSub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = bullet_text
        p_sub.font.size = Pt(16)
        p_sub.font.name = font_name
        p_sub.font.color.rgb = THEME["subtitle"]
        p_sub.alignment = PP_ALIGN.LEFT

    # Branding
    _add_branding(slide, language)


def _create_section_divider(prs: Presentation, section_num: int, title: str, total: int, language: str = "English"):
    """Create a section divider slide for major transitions."""
    font_name = _get_font(language)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, THEME["bg_alt"])
    _add_gradient_bar(slide, top=True)

    # Large section number
    txNum = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(2.0),
        width=Inches(2), height=Inches(1.5),
    )
    tf = txNum.text_frame
    p = tf.paragraphs[0]
    p.text = f"{section_num:02d}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    p.alignment = PP_ALIGN.LEFT

    # Section title
    txTitle = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(3.5),
        width=Inches(8), height=Inches(1.0),
    )
    tf_t = txTitle.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.name = font_name
    p_t.font.color.rgb = THEME["slide_title"]
    p_t.alignment = PP_ALIGN.LEFT

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(4.6),
        width=Inches(3), height=Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME["accent"]
    line.line.fill.background()

    _add_slide_number(slide, section_num, total, language)
    _add_branding(slide, language)


def _create_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    slide_num: int,
    total_slides: int,
    language: str = "English",
):
    """Create a content slide with title and bullet points."""
    font_name = _get_font(language)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Accent bar at top
    _add_gradient_bar(slide, top=True)

    # Small decorative circle
    _add_decorative_circle(slide, 8.5, 5.5, 1.5, RGBColor(0x12, 0x12, 0x2A))

    # Slide title
    txBox = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.4),
        width=Inches(8.4), height=Inches(0.9),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = THEME["slide_title"]
    p.alignment = PP_ALIGN.LEFT

    # Thin line under title
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.35),
        width=Inches(2.5), height=Inches(0.03),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = THEME["divider"]
    line_shape.line.fill.background()

    # Bullet points with alternating indent
    bullet_icons = ["▸", "◆", "●", "▹", "◇", "○"]
    txBody = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(1.7),
        width=Inches(8.4), height=Inches(5.0),
    )
    tf_body = txBody.text_frame
    tf_body.word_wrap = True

    for i, bullet in enumerate(bullets[:7]):  # Max 7 bullets
        if i == 0:
            p = tf_body.paragraphs[0]
        else:
            p = tf_body.add_paragraph()

        icon = bullet_icons[i % len(bullet_icons)]
        p.text = f"{icon}  {bullet}"
        p.font.size = Pt(16)
        p.font.name = font_name
        p.font.color.rgb = THEME["bullet"]
        p.space_after = Pt(14)
        p.alignment = PP_ALIGN.LEFT

    # Slide number + branding
    _add_slide_number(slide, slide_num, total_slides, language)
    _add_branding(slide, language)


def _create_references_slide(prs: Presentation, title: str, refs: list[str], slide_num: int, total: int, language: str = "English"):
    """Create a references slide with smaller text for URLs."""
    font_name = _get_font(language)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)
    _add_gradient_bar(slide, top=True)

    # Title
    txBox = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.4),
        width=Inches(8.4), height=Inches(0.8),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = THEME["slide_title"]

    # Divider
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.3),
        width=Inches(2.5), height=Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME["divider"]
    line.line.fill.background()

    # References
    txBody = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(1.6),
        width=Inches(8.4), height=Inches(5.0),
    )
    tf_body = txBody.text_frame
    tf_body.word_wrap = True

    for i, ref in enumerate(refs[:8]):
        if i == 0:
            p = tf_body.paragraphs[0]
        else:
            p = tf_body.add_paragraph()

        p.text = f"{i + 1}.  {ref}"
        p.font.size = Pt(12)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0xA0, 0xA0, 0xC0)
        p.space_after = Pt(10)

    _add_slide_number(slide, slide_num, total, language)
    _add_branding(slide, language)


def _create_thank_you_slide(prs: Presentation, topic: str, total: int, language: str = "English"):
    """Create a closing 'Thank You' slide."""
    font_name = _get_font(language)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, THEME["bg_alt"])
    _add_gradient_bar(slide, top=True, bottom=True)

    # Decorative
    _add_decorative_circle(slide, 1.0, 1.5, 3.0, RGBColor(0x15, 0x15, 0x30))
    _add_decorative_circle(slide, 6.0, 4.0, 2.5, RGBColor(0x18, 0x18, 0x35))

    # Thank You text
    txBox = slide.shapes.add_textbox(
        left=Inches(1.0), top=Inches(2.2),
        width=Inches(8.0), height=Inches(1.5),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = THEME["title"]
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txSub = slide.shapes.add_textbox(
        left=Inches(1.5), top=Inches(3.8),
        width=Inches(7.0), height=Inches(0.8),
    )
    tf_sub = txSub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = f"Research on: {topic}"
    p_sub.font.size = Pt(18)
    p_sub.font.name = font_name
    p_sub.font.color.rgb = THEME["subtitle"]
    p_sub.alignment = PP_ALIGN.CENTER

    # Branding
    txBrand = slide.shapes.add_textbox(
        left=Inches(2.5), top=Inches(5.0),
        width=Inches(5.0), height=Inches(0.5),
    )
    tf_b = txBrand.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "Powered by ScholarMind"
    p_b.font.size = Pt(14)
    p_b.font.name = font_name
    p_b.font.color.rgb = THEME["accent_light"]
    p_b.alignment = PP_ALIGN.CENTER


def generate_ppt(topic: str, slides_data: list[dict], language: str = "English") -> str:
    """
    Generate a professional PowerPoint presentation.

    Args:
        topic: The research topic.
        slides_data: List of dicts with 'title' and 'bullets' keys.
        language: Output language for font selection.

    Returns:
        The file path to the generated .pptx file.
    """
    filepath = os.path.join(OUTPUT_DIR, "presentation.pptx")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    total = len(slides_data) + 1  # +1 for thank you slide

    for i, slide_info in enumerate(slides_data):
        title = slide_info.get("title", f"Slide {i + 1}")
        bullets = slide_info.get("bullets", [])

        if i == 0:
            _create_title_slide(prs, title, bullets, language)
        elif i == len(slides_data) - 1:
            # Last data slide = references
            _create_references_slide(prs, title, bullets, i + 1, total, language)
        else:
            _create_content_slide(prs, title, bullets, i + 1, total, language)

    # Thank you slide
    _create_thank_you_slide(prs, topic, total, language)

    prs.save(filepath)
    print(f"[ScholarMind PPT] Generated: {filepath}")
    return filepath


if __name__ == "__main__":
    sample_slides = [
        {"title": "Sample Topic", "bullets": ["Research Presentation", "Powered by ScholarMind"]},
        {"title": "Executive Summary", "bullets": ["Key point 1", "Key point 2", "Key point 3", "Key point 4"]},
        {"title": "Background", "bullets": ["Historical context", "Origins and evolution", "Key milestones"]},
        {"title": "References", "bullets": ["https://example.com/paper1", "https://example.com/paper2"]},
    ]
    path = generate_ppt("Sample Topic", sample_slides)
    print(f"PPT saved to: {path}")
